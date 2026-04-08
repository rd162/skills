#!/usr/bin/env python3
"""
Document Converter Script
=========================
Converts documents (PDF, Word, PPT, Excel, draw.io diagrams, video) to:
1. Markdown using markitdown and docling libraries (dual versions)
2. WEBP images using pyvips (3-page sliding window for LLM vision processing)
3. draw.io diagrams: XML-parsed markdown + CLI-exported WEBP images
4. Video files: VTT subtitles via Whisper + scene-change cadre images via PySceneDetect
   - Manual VTT/SRT files from __SPECS__ are preserved alongside Whisper-generated VTTs

Also handles archives (ZIP, TAR, 7Z) — auto-extracts and processes all
documents inside using the same pipeline.

By default, scans the current working directory recursively for all
supported document types (excluding common non-content dirs like .git,
node_modules, __FRAGMENTS__, .venv, etc.).  If __SPECS__/ exists, it is
used as the scan root automatically.  Use --scan-dir to override.

Features:
- Recursive project scanning with smart directory exclusions
- Archive auto-extraction (ZIP, TAR, TAR.GZ, TGZ, 7Z)
- Incremental processing (SHA256 hash-based change detection)
- Dual markdown output (markitdown + docling) for cross-reference
- draw.io diagram support (XML parse + CLI image export)
- Sliding window WEBP images (pages 1-3, 2-4, 3-5 …)

Usage:
    # Scan project root recursively (auto-detects __SPECS__ if present):
    python scripts/doc_converter.py
    python scripts/doc_converter.py --force

    # Limit to a specific directory:
    python scripts/doc_converter.py --scan-dir __SPECS__
    python scripts/doc_converter.py --scan-dir docs/

    # Process a single file:
    python scripts/doc_converter.py --file "specific.pdf"

    # Backward-compatible (--specs-dir still works):
    python doc_converter.py --specs-dir __SPECS__ --fragments-dir __FRAGMENTS__

Requirements:
    pip install -r requirements_converter.txt
"""

import os
import sys
import json
import hashlib
import argparse
import tempfile
import shutil
import zipfile
import tarfile
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List
import logging

# ---------------------------------------------------------------------------
# vips dynamic module pre-loader
# ---------------------------------------------------------------------------
# On some macOS + libvips 8.18 installations, vips_init() does not call
# g_module_open for the vips-modules-X.Y directory, so dynamic modules such
# as vips-poppler (pdfload), vips-heif, vips-jxl are never registered.
# This function uses GLib's g_module_open directly — the same mechanism
# vips_init() is supposed to use — before pyvips is imported, so that by the
# time pyvips calls vips_init() the GObject types are already registered.
# It is a no-op when modules load correctly (fast path via _pdfload_available).

def _ensure_vips_modules_loaded() -> None:
    """Pre-load libvips dynamic modules via g_module_open if they are missing."""
    import ctypes, ctypes.util, os, glob as _glob

    # Locate libgmodule (GLib module system)
    gmodule_path = ctypes.util.find_library('gmodule-2.0')
    if not gmodule_path:
        # Common Homebrew location
        gmodule_path = '/opt/homebrew/opt/glib/lib/libgmodule-2.0.0.dylib'
    if not os.path.exists(gmodule_path or ''):
        return  # Can't find gmodule — skip silently

    try:
        gmodule = ctypes.cdll.LoadLibrary(gmodule_path)
        gmodule.g_module_open.restype = ctypes.c_void_p
        gmodule.g_module_open.argtypes = [ctypes.c_char_p, ctypes.c_int]
        gmodule.g_module_make_resident.argtypes = [ctypes.c_void_p]
    except Exception:
        return

    # Find the vips-modules directory (works for Homebrew and standard installs)
    candidates = [
        '/opt/homebrew/lib/vips-modules-*',
        '/usr/local/lib/vips-modules-*',
        '/usr/lib/vips-modules-*',
    ]
    module_dirs = []
    for pattern in candidates:
        module_dirs.extend(sorted(_glob.glob(pattern), reverse=True))  # newest first

    loaded = 0
    for module_dir in module_dirs:
        for dylib in sorted(Path(module_dir).glob('*.dylib')):
            handle = gmodule.g_module_open(str(dylib).encode(), 1)  # LAZY=1
            if handle:
                gmodule.g_module_make_resident(handle)  # prevent accidental unload
                loaded += 1
        if loaded:
            break  # found a working module dir

_vips_modules_loaded = False

def _lazy_vips_init():
    """Load vips modules on first use (not at import time) to avoid
    interfering with PyTorch/Whisper performance on ARM Macs."""
    global _vips_modules_loaded
    if not _vips_modules_loaded:
        _ensure_vips_modules_loaded()
        _vips_modules_loaded = True

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s | %(levelname)s | %(message)s',
    datefmt='%H:%M:%S'
)
logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Document types processed to markdown + WEBP
SUPPORTED_EXTENSIONS = {
    '.pdf':    'pdf',
    '.docx':   'word',
    '.doc':    'word',
    '.pptx':   'powerpoint',
    '.ppt':    'powerpoint',
    '.xlsx':   'excel',
    '.xls':    'excel',
    '.xlsm':   'excel',
    '.drawio': 'diagram',
    '.mp4':    'video',
    '.mkv':    'video',
    '.avi':    'video',
    '.mov':    'video',
    '.webm':   'video',
    '.m4v':    'video',
    '.wmv':    'video',
}

# Subtitle extensions that may accompany video files in __SPECS__
SUBTITLE_EXTENSIONS = {'.vtt', '.srt'}

# Default scene-detection threshold (lower = more sensitive, good for 30-min meetings)
VIDEO_SCENE_THRESHOLD = 5.0

# Default Whisper model for speech-to-text
VIDEO_WHISPER_MODEL = "base"

# Archive types — auto-extracted, contents processed recursively
ARCHIVE_EXTENSIONS = {
    '.zip',
    '.tar',
    '.tar.gz',
    '.tgz',
    '.7z',
}

# Directories to exclude when scanning a project root recursively.
# These are never scanned for source documents.
EXCLUDED_DIRS = {
    '__FRAGMENTS__',
    '.venv', 'venv', '.env',
    'node_modules',
    '.git', '.svn', '.hg',
    '.tmp', '.cache',
    '__pycache__', '.mypy_cache', '.pytest_cache',
    'dist', 'build', '.tox',
    '.idea', '.vscode',
}

IMAGE_DPI         = 150   # Resolution for PDF→image rendering
PAGES_PER_IMAGE   = 3     # Sliding window width (pages per WEBP)
WEBP_QUALITY      = 85    # WEBP encoder quality (0–100)
MAX_ARCHIVE_DEPTH = 2     # Prevent infinite nesting

# ---------------------------------------------------------------------------
# draw.io CLI detection (for .drawio → PDF/PNG export)
# ---------------------------------------------------------------------------

def _detect_drawio_cli() -> Optional[str]:
    """Find draw.io CLI executable. Returns path or None."""
    candidates = [
        # macOS .app bundle
        "/Applications/draw.io.app/Contents/MacOS/draw.io",
        os.path.expanduser("~/Applications/draw.io.app/Contents/MacOS/draw.io"),
        # Linux / Homebrew / snap
        "drawio",
        "draw.io",
        "/opt/homebrew/bin/drawio",
        "/usr/local/bin/drawio",
        "/snap/bin/drawio",
    ]
    for candidate in candidates:
        if os.path.isfile(candidate):
            return candidate
    # Check PATH for 'drawio' or 'draw.io'
    for name in ("drawio", "draw.io"):
        result = shutil.which(name)
        if result:
            return result
    return None

DRAWIO_CLI = _detect_drawio_cli()

# ---------------------------------------------------------------------------
# Manifest helpers
# ---------------------------------------------------------------------------

def get_file_hash(filepath: Path) -> str:
    """SHA256 hash of a file (change-detection key)."""
    sha256 = hashlib.sha256()
    with open(filepath, "rb") as fh:
        for block in iter(lambda: fh.read(4096), b""):
            sha256.update(block)
    return sha256.hexdigest()


def load_manifest(manifest_path: Path) -> Dict:
    if manifest_path.exists():
        with open(manifest_path, 'r') as fh:
            return json.load(fh)
    return {"files": {}, "last_updated": None}


def save_manifest(manifest_path: Path, manifest: Dict):
    manifest["last_updated"] = datetime.now().isoformat()
    with open(manifest_path, 'w') as fh:
        json.dump(manifest, fh, indent=2)


def _manifest_key(filepath: Path) -> str:
    """Stable manifest key: path relative to CWD (the project root).

    Always resolves relative to the current working directory, regardless
    of which --scan-dir was used.  This ensures the same file always gets
    the same key whether you run with --scan-dir docs, --scan-dir .,
    or no flag at all.

    Using relative paths avoids key collisions when files with the same
    name live in different subdirectories (e.g. docs/req.pdf vs specs/req.pdf).
    """
    cwd = Path.cwd().resolve()
    try:
        return str(filepath.resolve().relative_to(cwd))
    except ValueError:
        # Fallback for files outside CWD (e.g. archive-extracted temps)
        return str(filepath.name)


def needs_processing(filepath: Path, manifest: Dict, force: bool = False,
                     file_key: Optional[str] = None) -> bool:
    """Check whether a file needs (re-)processing.

    Compares the file's current SHA256 hash against the manifest record.
    Returns True if:  force flag is set, file is new, or hash differs.
    """
    if force:
        return True
    key = file_key or filepath.name
    current_hash = get_file_hash(filepath)
    if key not in manifest["files"]:
        return True
    prev = manifest["files"][key]
    if prev.get("hash") != current_hash:
        return True
    # Reprocess if previous run failed or produced no images (upgrade markdown_only)
    if prev.get("status") in ("error", "failed", "pending", "markdown_only"):
        return True
    return False


def _detect_orphaned_entries(manifest: Dict, current_keys: set) -> List[str]:
    """Find manifest entries whose source documents are no longer present.

    An entry is orphaned when:
      - Its key is not in the set of currently-scanned file keys, AND
      - Its recorded source path no longer exists on disk.

    Archive-internal entries (type == 'archive' children) are considered
    orphaned when their parent archive is orphaned.
    """
    orphaned = []
    for key, info in list(manifest.get("files", {}).items()):
        if key in current_keys:
            continue
        # Double-check: the source file might still exist even if it wasn't
        # in the scan set (e.g. user changed --scan-dir between runs)
        source = info.get("source", "")
        if source and Path(source).exists():
            continue
        orphaned.append(key)
    return sorted(orphaned)


def _clean_orphaned_fragments(
    orphaned_keys: List[str], manifest: Dict, fragments_dir: Path
) -> int:
    """Remove fragment directories and manifest entries for orphaned docs.

    Returns the number of entries removed.
    """
    cleaned = 0
    for key in orphaned_keys:
        info = manifest["files"].get(key, {})

        # Determine the fragment directory name (mirrors process_document logic)
        display_name = info.get("source", key)
        safe_stem = Path(display_name).stem
        frag_dir = fragments_dir / safe_stem

        if frag_dir.exists() and frag_dir.is_dir():
            shutil.rmtree(frag_dir)
            logger.info(f"  🗑 Cleaned orphaned fragments: {frag_dir.name}/ (was: {key})")

        # Remove archive children if this was an archive
        for child_key in info.get("archive_contents", []):
            if child_key in manifest["files"]:
                child_stem = Path(child_key).stem
                child_dir = fragments_dir / child_stem
                if child_dir.exists() and child_dir.is_dir():
                    shutil.rmtree(child_dir)
                    logger.info(f"  🗑 Cleaned orphaned archive content: {child_dir.name}/")
                del manifest["files"][child_key]

        del manifest["files"][key]
        cleaned += 1
    return cleaned


# ---------------------------------------------------------------------------
# Archive extraction
# ---------------------------------------------------------------------------

def is_archive(filepath: Path) -> bool:
    """Return True if the file is a supported archive."""
    name = filepath.name.lower()
    for ext in ARCHIVE_EXTENSIONS:
        if name.endswith(ext):
            return True
    return False


def extract_archive(filepath: Path, dest_dir: Path) -> List[Path]:
    """
    Extract a supported archive into dest_dir.
    Returns a list of extracted file paths (non-archive documents only).
    Nested archives are extracted recursively up to MAX_ARCHIVE_DEPTH.
    """
    extracted: List[Path] = []
    name_lower = filepath.name.lower()

    try:
        if name_lower.endswith('.zip'):
            with zipfile.ZipFile(filepath, 'r') as zf:
                zf.extractall(dest_dir)

        elif name_lower.endswith(('.tar.gz', '.tgz')):
            with tarfile.open(filepath, 'r:gz') as tf:
                tf.extractall(dest_dir)

        elif name_lower.endswith('.tar'):
            with tarfile.open(filepath, 'r:*') as tf:
                tf.extractall(dest_dir)

        elif name_lower.endswith('.7z'):
            try:
                import py7zr
                with py7zr.SevenZipFile(filepath, mode='r') as szf:
                    szf.extractall(path=dest_dir)
            except ImportError:
                logger.warning("⚠ py7zr not installed — skipping .7z archive. "
                               "Install with: pip install py7zr")
                return []
        else:
            logger.warning(f"⚠ Unsupported archive format: {filepath.name}")
            return []

    except Exception as exc:
        logger.error(f"✗ Archive extraction failed for {filepath.name}: {exc}")
        return []

    # Walk extracted tree and collect processable files
    for candidate in sorted(dest_dir.rglob("*")):
        if not candidate.is_file():
            continue
        if candidate.suffix.lower() in SUPPORTED_EXTENSIONS:
            extracted.append(candidate)
        # One level of nested archive support (depth guard handled by caller)
        elif is_archive(candidate):
            nested_dir = dest_dir / (candidate.stem + "_extracted")
            nested_dir.mkdir(parents=True, exist_ok=True)
            extracted.extend(extract_archive(candidate, nested_dir))

    logger.info(f"  ✓ Archive extracted: {len(extracted)} document(s) found")
    return extracted

# ---------------------------------------------------------------------------
# Markdown converters
# ---------------------------------------------------------------------------

def convert_with_markitdown(filepath: Path, output_dir: Path) -> Optional[Path]:
    """Convert document to markdown using markitdown."""
    # B4: Legacy .doc (Word 97-2003) files cannot be reliably converted by markitdown
    # (markitdown misidentifies them as Excel files via OLE2 compound document).
    # Skip gracefully — docling and WEBP fallback cover these files.
    if filepath.suffix.lower() == '.doc':
        logger.warning(
            f"  ⚠ Skipping markitdown for legacy .doc file: {filepath.name} "
            f"(Word 97-2003 format not supported — use LibreOffice to convert to .docx first)"
        )
        return None

    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(str(filepath))
        output_file = output_dir / f"{filepath.stem}_markitdown.md"
        with open(output_file, 'w', encoding='utf-8') as fh:
            fh.write(f"<!-- Source: {filepath.name} -->\n")
            fh.write(f"<!-- Converter: markitdown -->\n")
            fh.write(f"<!-- Generated: {datetime.now().isoformat()} -->\n\n")
            fh.write(result.text_content)
        logger.info(f"  ✓ markitdown → {output_file.name}")
        return output_file
    except ImportError:
        logger.error("  ✗ markitdown not installed: pip install 'markitdown[all]'")
        return None
    except Exception as exc:
        logger.error(f"  ✗ markitdown failed for {filepath.name}: {exc}")
        return None
    except BaseException as exc:
        # B1: Python 3.14 compatibility — markitdown.FileConversionException and similar
        # custom exceptions sometimes inherit from BaseException directly in newer versions,
        # bypassing a plain `except Exception` handler.  Catch here and log gracefully.
        logger.warning(
            f"  ✗ markitdown raised non-standard exception for {filepath.name}: "
            f"{type(exc).__name__}: {exc}"
        )
        return None


def convert_with_docling(filepath: Path, output_dir: Path) -> Optional[Path]:
    """Convert document to markdown using docling.

    Runs docling in an isolated subprocess so that a crash (segfault, OOM)
    in docling's OCR/ML pipeline does not kill the main converter process.
    """
    import subprocess, sys
    output_file = output_dir / f"{filepath.stem}_docling.md"

    script = f"""
import sys
from pathlib import Path
from datetime import datetime
from docling.document_converter import DocumentConverter

filepath = Path({str(filepath)!r})
output_file = Path({str(output_file)!r})
converter = DocumentConverter()
result = converter.convert(str(filepath))
markdown_content = result.document.export_to_markdown()
with open(output_file, 'w', encoding='utf-8') as fh:
    fh.write(f"<!-- Source: {{filepath.name}} -->\\n")
    fh.write(f"<!-- Converter: docling -->\\n")
    fh.write(f"<!-- Generated: {{datetime.now().isoformat()}} -->\\n\\n")
    fh.write(markdown_content)
"""
    try:
        r = subprocess.run(
            [sys.executable, '-c', script],
            capture_output=True, text=True, timeout=300,
        )
        if r.returncode == 0 and output_file.exists():
            logger.info(f"  ✓ docling → {output_file.name}")
            return output_file
        else:
            err = (r.stderr or '').strip().splitlines()
            short = err[-1] if err else f'exit {r.returncode}'
            logger.error(f"  ✗ docling failed for {filepath.name}: {short}")
            return None
    except subprocess.TimeoutExpired:
        logger.error(f"  ✗ docling timed out for {filepath.name}")
        return None
    except Exception as exc:
        logger.error(f"  ✗ docling error for {filepath.name}: {exc}")
        return None


# ---------------------------------------------------------------------------
# draw.io helpers
# ---------------------------------------------------------------------------

def _parse_drawio_xml(filepath: Path) -> str:
    """
    Parse a .drawio XML file and extract all text content, node labels,
    edge labels, and diagram metadata into a readable markdown document.

    .drawio files are XML (mxGraphModel). Text lives in:
      - mxCell 'value' attributes (node/edge labels)
      - <mxCell> child text nodes (rare)
      - diagram 'name' attribute (tab/page names)
    """
    import xml.etree.ElementTree as ET

    try:
        tree = ET.parse(filepath)
        root = tree.getroot()
    except ET.ParseError as exc:
        logger.warning(f"  ⚠ draw.io XML parse error: {exc}")
        return f"# {filepath.stem}\n\n*Failed to parse draw.io XML: {exc}*\n"

    lines = [f"# {filepath.stem}\n"]

    # Iterate over each diagram tab (multi-page drawio support)
    diagrams = root.findall('.//diagram')
    if not diagrams:
        # Might be a flat mxGraphModel without <diagram> wrapper
        diagrams = [root]

    for diagram_idx, diagram in enumerate(diagrams):
        diagram_name = diagram.get('name', f'Page {diagram_idx + 1}')
        lines.append(f"\n## {diagram_name}\n")

        # Collect nodes and edges in two passes:
        #   Pass 1: identify edge cells (have edge="1", source/target, or endArrow/edgeStyle)
        #   Pass 2: classify remaining cells as nodes or edge-labels
        #
        # Edge labels in drawio are separate mxCell elements with:
        #   - connectable="0"  and/or  style containing "edgeLabel"
        #   - parent attribute pointing to the edge cell's id
        # They must NOT be classified as nodes.

        import re

        nodes = []
        edges = {}       # edge_id → {label, source, target}
        edge_labels = {} # edge_id → label text (from child edgeLabel cells)
        all_cells = []   # raw cell data for two-pass processing

        # ── Pass 1: gather raw cell data and identify edges ──
        for cell in diagram.iter('mxCell'):
            cell_id = cell.get('id', '')
            value = (cell.get('value') or '').strip()
            style = cell.get('style', '')
            source = cell.get('source')
            target = cell.get('target')
            parent = cell.get('parent', '')
            is_edge_attr = cell.get('edge', '') == '1'
            connectable = cell.get('connectable', '')

            # Skip structural root cells (id 0 and 1 with no content)
            if cell_id in ('0', '1') and not value:
                continue

            # Strip HTML tags from labels (drawio uses <br>, <b>, <div>, etc.)
            clean_value = ''
            if value:
                clean_value = re.sub(r'<[^>]+>', ' ', value).strip()
                clean_value = re.sub(r'\s+', ' ', clean_value)

            # Detect edge cells: explicit edge="1" attr, or has source/target,
            # or style contains endArrow/edgeStyle keywords
            is_edge = (
                is_edge_attr
                or (source and target)
                or 'endArrow' in style
                or 'edgeStyle' in style
            )

            # Detect edge-label cells: connectable="0" or style has "edgeLabel"
            is_edge_label = (
                connectable == '0'
                or 'edgeLabel' in style
            )

            all_cells.append({
                'id': cell_id,
                'value': clean_value,
                'style': style,
                'source': source,
                'target': target,
                'parent': parent,
                'is_edge': is_edge,
                'is_edge_label': is_edge_label,
            })

            if is_edge:
                edges[cell_id] = {
                    'label': clean_value,
                    'source': source,
                    'target': target,
                }

        # ── Pass 2: attach edge labels, classify remaining as nodes ──
        for cell in all_cells:
            if cell['is_edge']:
                continue  # already in edges dict

            if cell['is_edge_label'] and cell['parent'] in edges:
                # This is a label attached to a known edge — merge it
                if cell['value']:
                    edges[cell['parent']]['label'] = cell['value']
                continue

            if not cell['value']:
                continue

            # Detect shape type from style for node classification
            style = cell['style']
            shape = 'component'
            if 'cylinder' in style:
                shape = 'database'
            elif 'rhombus' in style or 'diamond' in style:
                shape = 'decision'
            elif 'ellipse' in style:
                shape = 'actor'
            elif 'cloud' in style:
                shape = 'cloud'
            elif 'shape=mxgraph.aws' in style or 'shape=mxgraph.azure' in style:
                shape = 'cloud-service'

            nodes.append({
                'id': cell['id'],
                'label': cell['value'],
                'shape': shape,
            })

        if nodes:
            lines.append("\n### Components\n")
            for node in nodes:
                shape_hint = f" ({node['shape']})" if node['shape'] != 'component' else ''
                lines.append(f"- **{node['label']}**{shape_hint}")

        edge_list = list(edges.values())
        if edge_list:
            lines.append("\n### Connections\n")
            # Build an ID→label map for readable edge descriptions
            id_to_label = {n['id']: n['label'] for n in nodes}
            for edge in edge_list:
                src_label = id_to_label.get(edge['source'], edge['source'] or '?')
                tgt_label = id_to_label.get(edge['target'], edge['target'] or '?')
                if edge['label']:
                    lines.append(f"- {src_label} → {tgt_label} ({edge['label']})")
                else:
                    lines.append(f"- {src_label} → {tgt_label}")

    return '\n'.join(lines) + '\n'


def convert_drawio_to_markdown(filepath: Path, output_dir: Path) -> Optional[Path]:
    """Convert .drawio file to markdown by parsing its XML structure."""
    try:
        markdown_content = _parse_drawio_xml(filepath)
        output_file = output_dir / f"{filepath.stem}_drawio_parsed.md"
        with open(output_file, 'w', encoding='utf-8') as fh:
            fh.write(f"<!-- Source: {filepath.name} -->\n")
            fh.write(f"<!-- Converter: drawio-xml-parser -->\n")
            fh.write(f"<!-- Generated: {datetime.now().isoformat()} -->\n\n")
            fh.write(markdown_content)
        logger.info(f"  ✓ drawio-xml-parser → {output_file.name}")
        return output_file
    except Exception as exc:
        logger.error(f"  ✗ drawio XML parse failed for {filepath.name}: {exc}")
        return None


def convert_drawio_to_images(
    filepath: Path, output_dir: Path, temp_dir: Path
) -> List[Path]:
    """
    Convert .drawio to WEBP images.

    Strategy:
      1. draw.io CLI → PDF → pyvips WEBP (best quality)
      2. draw.io CLI → PNG → WEBP (fallback)
      3. No CLI → skip images, rely on XML-parsed markdown
    """
    if not DRAWIO_CLI:
        logger.warning(
            "  ⚠ draw.io CLI not found — skipping image export. "
            "Install draw.io desktop app for diagram image rendering."
        )
        return []

    import subprocess
    base_name = filepath.stem
    images: List[Path] = []

    # Count pages (diagrams) in the file
    import xml.etree.ElementTree as ET
    try:
        tree = ET.parse(filepath)
        diagram_count = len(tree.getroot().findall('.//diagram'))
        if diagram_count == 0:
            diagram_count = 1
    except Exception:
        diagram_count = 1

    # Export each page
    for page_idx in range(diagram_count):
        # Try PDF first (feeds into existing WEBP pipeline)
        pdf_path = temp_dir / f"{base_name}_p{page_idx}.pdf"
        try:
            cmd = [
                DRAWIO_CLI, '--export',
                '--format', 'pdf',
                '--page-index', str(page_idx),
                '--output', str(pdf_path),
                str(filepath),
            ]
            result = subprocess.run(
                cmd, capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0 and pdf_path.exists() and pdf_path.stat().st_size > 0:
                # Convert PDF page to single WEBP
                page_images = render_pdf_pages_to_images(
                    pdf_path, output_dir, f"{base_name}_diagram"
                )
                images.extend(page_images)
                continue
        except subprocess.TimeoutExpired:
            logger.warning(f"  ⚠ draw.io PDF export timed out for page {page_idx}")
        except Exception as exc:
            logger.warning(f"  ⚠ draw.io PDF export failed for page {page_idx}: {exc}")

        # Fallback: try PNG export
        png_path = temp_dir / f"{base_name}_p{page_idx}.png"
        try:
            cmd = [
                DRAWIO_CLI, '--export',
                '--format', 'png',
                '--scale', '2',
                '--page-index', str(page_idx),
                '--output', str(png_path),
                str(filepath),
            ]
            result = subprocess.run(
                cmd, capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0 and png_path.exists() and png_path.stat().st_size > 0:
                # Convert PNG to WEBP
                try:
                    _lazy_vips_init(); import pyvips
                    vips_img = pyvips.Image.new_from_file(str(png_path))
                    webp_name = f"{base_name}_diagram_p{page_idx + 1:03d}.webp"
                    webp_path = output_dir / webp_name
                    vips_img.webpsave(str(webp_path), Q=WEBP_QUALITY)
                    images.append(webp_path)
                    logger.info(f"  ✓ drawio page {page_idx + 1} → {webp_name}")
                    continue
                except Exception as vips_exc:
                    # Last resort: keep the PNG as-is (still useful for vision)
                    from shutil import copy2
                    kept_png = output_dir / f"{base_name}_diagram_p{page_idx + 1:03d}.png"
                    copy2(png_path, kept_png)
                    images.append(kept_png)
                    logger.info(f"  ✓ drawio page {page_idx + 1} → {kept_png.name} (PNG, no vips)")
                    continue
        except subprocess.TimeoutExpired:
            logger.warning(f"  ⚠ draw.io PNG export timed out for page {page_idx}")
        except Exception as exc:
            logger.warning(f"  ⚠ draw.io PNG export failed for page {page_idx}: {exc}")

        logger.warning(f"  ✗ Could not export page {page_idx} of {filepath.name}")

    return images


# ---------------------------------------------------------------------------
# PDF helpers
# ---------------------------------------------------------------------------

def get_pdf_page_count(filepath: Path) -> int:
    """Return the number of pages in a PDF."""
    try:
        _lazy_vips_init(); import pyvips
        img = pyvips.Image.pdfload(str(filepath), n=1)
        return img.get('n-pages')
    except Exception:
        pass
    try:
        import fitz
        doc = fitz.open(str(filepath))
        count = doc.page_count
        doc.close()
        return count
    except Exception:
        pass
    try:
        import subprocess
        r = subprocess.run(
            ['pdfinfo', str(filepath)], capture_output=True, text=True, timeout=30
        )
        for line in r.stdout.splitlines():
            if line.lower().startswith('pages:'):
                return int(line.split(':')[1].strip())
    except Exception:
        pass
    return 1

# ---------------------------------------------------------------------------
# Office → PDF converters (produce PDF for render_pdf_pages_to_images)
# ---------------------------------------------------------------------------

# Candidate Chrome/Chromium executable paths (in priority order)
_CHROME_CANDIDATES = [
    "google-chrome",
    "google-chrome-stable",
    "chromium",
    "chromium-browser",
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    "/Applications/Chromium.app/Contents/MacOS/Chromium",
    "/usr/bin/google-chrome",
    "/usr/bin/chromium",
]


def _find_chrome() -> Optional[str]:
    """Return the first usable Chrome/Chromium executable, or None."""
    import shutil
    for candidate in _CHROME_CANDIDATES:
        if Path(candidate).is_file():
            return candidate
        found = shutil.which(candidate)
        if found:
            return found
    return None


def _office_to_pdf_libreoffice(filepath: Path, temp_dir: Path) -> Optional[Path]:
    """Convert Office document to PDF via LibreOffice headless."""
    import subprocess
    output_pdf = temp_dir / f"{filepath.stem}.pdf"
    cmd = [
        'soffice', '--headless', '--convert-to', 'pdf',
        '--outdir', str(temp_dir), str(filepath)
    ]
    try:
        subprocess.run(cmd, capture_output=True, text=True, timeout=120, check=False)
        return output_pdf if output_pdf.exists() else None
    except FileNotFoundError:
        logger.debug("  LibreOffice (soffice) not found")
        return None
    except Exception as exc:
        logger.debug(f"  LibreOffice conversion error: {exc}")
        return None


def _office_to_pdf_chrome(filepath: Path, temp_dir: Path) -> Optional[Path]:
    """Convert DOCX to PDF via mammoth (DOCX→HTML) + Chrome headless (HTML→PDF).

    Fully automated — no GUI, no dialogs, no AppleScript.
    Preserves embedded images (base64 in HTML), tables, and text formatting.
    Requires: mammoth (pip) + Google Chrome or Chromium (system).
    """
    import subprocess
    ext = filepath.suffix.lower()
    if ext not in ('.docx', '.doc'):
        return None

    chrome = _find_chrome()
    if not chrome:
        logger.debug("  Chrome/Chromium not found — skipping mammoth+Chrome strategy")
        return None

    try:
        import mammoth
    except ImportError:
        logger.debug("  mammoth not installed — skipping mammoth+Chrome strategy")
        return None

    try:
        # Step 1: DOCX → HTML via mammoth (embedded images become base64 data URIs)
        html_path = temp_dir / f"{filepath.stem}.html"
        with open(filepath, "rb") as f:
            result = mammoth.convert_to_html(f)
        html_content = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
  body {{ font-family: Calibri, Arial, sans-serif; font-size: 11pt;
          margin: 2cm; line-height: 1.4; }}
  table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
  td, th {{ border: 1px solid #ccc; padding: 4px 8px; vertical-align: top; }}
  img {{ max-width: 100%; height: auto; display: block; margin: 1em 0; }}
  h1 {{ font-size: 18pt; }} h2 {{ font-size: 14pt; }} h3 {{ font-size: 12pt; }}
  p {{ margin: 0.5em 0; }}
</style>
</head>
<body>
{result.value}
</body>
</html>"""
        html_path.write_text(html_content, encoding="utf-8")

        # Step 2: HTML → PDF via Chrome headless (preserves images and layout)
        output_pdf = temp_dir / f"{filepath.stem}.pdf"
        cmd = [
            chrome,
            "--headless=new",
            "--no-sandbox",
            "--disable-gpu",
            "--disable-software-rasterizer",
            "--disable-dev-shm-usage",
            f"--print-to-pdf={output_pdf}",
            "--print-to-pdf-no-header",
            f"file://{html_path.resolve()}",
        ]
        subprocess.run(cmd, capture_output=True, text=True, timeout=60, check=False)
        return output_pdf if output_pdf.exists() else None

    except Exception as exc:
        logger.debug(f"  mammoth+Chrome conversion error: {exc}")
        return None


def _office_to_pdf_docx2pdf(filepath: Path, temp_dir: Path) -> Optional[Path]:
    """Convert DOCX to PDF via docx2pdf (uses Word on macOS/Windows).

    Note: on macOS this drives Word via AppleScript and may show permission
    dialogs or time out if Word has not been granted Automation access.
    Kept as a last-resort fallback after LibreOffice and Chrome strategies.
    """
    try:
        from docx2pdf import convert as d2p_convert
        output_pdf = temp_dir / f"{filepath.stem}.pdf"
        d2p_convert(str(filepath), str(output_pdf))
        return output_pdf if output_pdf.exists() else None
    except ImportError:
        logger.debug("  docx2pdf not installed")
        return None
    except Exception as exc:
        logger.debug(f"  docx2pdf conversion error: {exc}")
        return None


def convert_office_to_pdf(filepath: Path, temp_dir: Path) -> Optional[Path]:
    """
    Try multiple strategies to produce a PDF from an Office document.

    Order for DOCX/DOC:
      1. LibreOffice headless       (cross-platform, best fidelity)
      2. mammoth + Chrome headless  (no install needed if Chrome present;
                                     fully automated, preserves images)
      3. docx2pdf                   (macOS Word / Windows — last resort;
                                     may show permission dialogs)

    Order for PPTX/PPT:
      1. LibreOffice headless
    """
    ext = filepath.suffix.lower()

    # --- Strategy 1: LibreOffice (works for all Office types) ---
    logger.info(f"  Trying LibreOffice → PDF…")
    pdf = _office_to_pdf_libreoffice(filepath, temp_dir)
    if pdf:
        logger.info(f"  ✓ LibreOffice → {pdf.name}")
        return pdf

    # --- Strategy 2: mammoth + Chrome headless (DOCX/DOC only) ---
    if ext in ('.docx', '.doc'):
        logger.info(f"  Trying mammoth + Chrome headless → PDF…")
        pdf = _office_to_pdf_chrome(filepath, temp_dir)
        if pdf:
            logger.info(f"  ✓ mammoth + Chrome → {pdf.name}")
            return pdf

    # --- Strategy 3: docx2pdf / Word (DOCX/DOC only, last resort) ---
    if ext in ('.docx', '.doc'):
        logger.info(f"  Trying docx2pdf → PDF…")
        pdf = _office_to_pdf_docx2pdf(filepath, temp_dir)
        if pdf:
            logger.info(f"  ✓ docx2pdf → {pdf.name}")
            return pdf

    logger.warning(f"  ⚠ Could not convert {filepath.name} to PDF via any converter")
    return None

# ---------------------------------------------------------------------------
# DOCX native fallback (no external tool required)
# ---------------------------------------------------------------------------

def convert_docx_to_images_fallback(
    filepath: Path,
    output_dir: Path,
    base_name: str,
) -> List[Path]:
    """
    Fallback DOCX → WEBP when no PDF converter is available.

    Strategy:
      - Extract all embedded images from the DOCX (in document order)
      - Group them in sliding windows identical to the PDF renderer
      - If no images are embedded, render text paragraphs as plain-text tiles

    This produces lower-fidelity output than PDF rendering but preserves
    all visual content embedded in the Word document.
    """
    try:
        from docx import Document as DocxDocument
        from PIL import Image, ImageDraw, ImageFont
        import io
    except ImportError as exc:
        logger.warning(f"  ⚠ DOCX fallback unavailable ({exc})")
        return []

    try:
        doc = DocxDocument(str(filepath))

        # --- Collect page-like units: embedded images first, then text blocks ---
        units: List[Image.Image] = []

        # 1. Embedded images (in-line shapes)
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    img_bytes = rel.target_part.blob
                    pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    units.append(pil_img)
                except Exception:
                    pass

        # 2. If no images, render text paragraphs as pseudo-pages
        if not units:
            TEXT_W, TEXT_H = 1240, 1754   # A4-ish at ~150 dpi
            FONT_SIZE      = 18
            LINE_H         = FONT_SIZE + 6
            MARGIN         = 60
            MAX_LINES      = (TEXT_H - 2 * MARGIN) // LINE_H

            try:
                font = ImageFont.truetype("DejaVuSans.ttf", FONT_SIZE)
            except Exception:
                font = ImageFont.load_default()

            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            # Chunk paragraphs into pseudo-pages
            for i in range(0, max(len(paragraphs), 1), MAX_LINES):
                chunk = paragraphs[i: i + MAX_LINES]
                page_img = Image.new("RGB", (TEXT_W, TEXT_H), color=(255, 255, 255))
                draw = ImageDraw.Draw(page_img)
                y = MARGIN
                for line in chunk:
                    # Wrap long lines
                    max_chars = (TEXT_W - 2 * MARGIN) // (FONT_SIZE // 2)
                    while len(line) > max_chars:
                        draw.text((MARGIN, y), line[:max_chars], fill=(30, 30, 30), font=font)
                        line = line[max_chars:]
                        y += LINE_H
                        if y > TEXT_H - MARGIN:
                            break
                    if y <= TEXT_H - MARGIN:
                        draw.text((MARGIN, y), line, fill=(30, 30, 30), font=font)
                        y += LINE_H
                units.append(page_img)

        if not units:
            logger.warning(f"  ⚠ DOCX fallback: no content extracted from {filepath.name}")
            return []

        # --- Apply same sliding window as render_pdf_pages_to_images ---
        total = len(units)
        logger.info(f"  DOCX fallback: {total} unit(s) → sliding window")

        created: List[Path] = []
        for start in range(total):
            end = min(start + PAGES_PER_IMAGE, total)
            window = units[start:end]

            # Normalise widths for vertical stack
            target_w = max(img.width for img in window)
            normalised = []
            for img in window:
                if img.width != target_w:
                    scale = target_w / img.width
                    new_h = int(img.height * scale)
                    img = img.resize((target_w, new_h), Image.LANCZOS)
                normalised.append(img)

            total_h = sum(img.height for img in normalised)
            composite = Image.new("RGB", (target_w, total_h), (255, 255, 255))
            y_off = 0
            for img in normalised:
                composite.paste(img, (0, y_off))
                y_off += img.height

            out_file = output_dir / f"{base_name}_p{start+1:03d}-{end:03d}.webp"
            composite.save(str(out_file), "WEBP", quality=WEBP_QUALITY)
            created.append(out_file)
            logger.info(f"  ✓ DOCX fallback image: units {start+1}-{end} → {out_file.name}")

            if end >= total:
                break

        return created

    except Exception as exc:
        logger.warning(f"  ⚠ DOCX native fallback failed: {exc}")
        return []

# ---------------------------------------------------------------------------
# PPTX native fallback
# ---------------------------------------------------------------------------

def convert_pptx_to_images_fallback(
    filepath: Path,
    output_dir: Path,
    base_name: str,
) -> List[Path]:
    """Fallback: extract slide images from PPTX using python-pptx + PIL."""
    try:
        from pptx import Presentation
        from PIL import Image
        import io
    except ImportError as exc:
        logger.warning(f"  ⚠ PPTX fallback unavailable ({exc})")
        return []

    try:
        prs = Presentation(str(filepath))
        images: List[Path] = []

        for idx, slide in enumerate(prs.slides):
            canvas = Image.new('RGB', (1920, 1080), color='white')
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        slide_img = Image.open(io.BytesIO(shape.image.blob))
                        canvas.paste(slide_img, (0, 0))
                        break
                    except Exception:
                        pass
            img_path = output_dir / f"{base_name}_p{idx+1:03d}-{idx+1:03d}.webp"
            canvas.save(str(img_path), 'WEBP', quality=WEBP_QUALITY)
            images.append(img_path)

        return images
    except Exception as exc:
        logger.warning(f"  ⚠ PPTX fallback failed: {exc}")
        return []

# ---------------------------------------------------------------------------
# Core PDF page renderer (pyvips sliding window)
# ---------------------------------------------------------------------------

def _pdfload_available() -> bool:
    """Return True if pyvips has a working pdfload operation."""
    try:
        _lazy_vips_init(); import pyvips
        # Probe with a known-bad path — if the op exists we get a file error,
        # if the op is missing we get AttributeError before any I/O.
        pyvips.Image.pdfload("/dev/null", page=0, n=1)
    except AttributeError:
        return False
    except Exception:
        return True  # op exists, file error is expected
    return True


def _render_pdf_via_pdftoppm(
    pdf_path: Path,
    output_dir: Path,
    base_name: str,
    pages_per_image: int = PAGES_PER_IMAGE,
) -> List[Path]:
    """
    Fallback renderer: pdftoppm → PNG files → pyvips WEBP sliding window.

    Used when pyvips pdfload (poppler dynamic module) is unavailable.
    Requires: pdftoppm (poppler-utils) + pyvips (for PNG→WEBP conversion).
    """
    import subprocess
    import shutil

    if not shutil.which('pdftoppm'):
        logger.error("  ✗ pdftoppm not found — install poppler (brew install poppler / apt install poppler-utils)")
        return []

    try:
        _lazy_vips_init(); import pyvips
    except ImportError:
        logger.error("  ✗ pyvips not installed")
        return []

    import tempfile
    created: List[Path] = []

    with tempfile.TemporaryDirectory() as tmp_str:
        tmp = Path(tmp_str)
        prefix = tmp / "page"

        r = subprocess.run(
            ['pdftoppm', '-r', str(IMAGE_DPI), '-png', str(pdf_path), str(prefix)],
            capture_output=True, text=True, timeout=120,
        )
        if r.returncode != 0:
            logger.error(f"  ✗ pdftoppm failed: {r.stderr[:300]}")
            return []

        page_files = sorted(tmp.glob("page-*.png"))
        if not page_files:
            page_files = sorted(tmp.glob("page*.png"))  # some versions omit the dash

        total_pages = len(page_files)
        if total_pages == 0:
            logger.warning("  ⚠ pdftoppm produced no pages")
            return []

        logger.info(f"  pdftoppm: {total_pages} page(s) rendered")

        for start_page in range(total_pages):
            end_page = min(start_page + pages_per_image, total_pages)
            page_imgs = []

            for page_num in range(start_page, end_page):
                vimg = pyvips.Image.new_from_file(
                    str(page_files[page_num]), access='sequential'
                )
                if vimg.bands == 4:
                    vimg = vimg.flatten(background=[255, 255, 255])
                page_imgs.append(vimg)

            img = (
                pyvips.Image.arrayjoin(page_imgs, across=1)
                if len(page_imgs) > 1
                else page_imgs[0]
            )

            out_file = output_dir / f"{base_name}_p{start_page+1:03d}-{end_page:03d}.webp"
            img.webpsave(str(out_file), Q=WEBP_QUALITY, effort=4, smart_subsample=True)
            created.append(out_file)
            logger.info(f"  ✓ Image: pages {start_page+1}-{end_page} → {out_file.name}")

            if end_page >= total_pages:
                break

    return created


def render_pdf_pages_to_images(
    pdf_path: Path,
    output_dir: Path,
    base_name: str,
    pages_per_image: int = PAGES_PER_IMAGE,
) -> List[Path]:
    """
    Render PDF pages to WEBP images with a sliding window.

    Strategy:
      1. pyvips pdfload  — fastest; requires poppler dynamic module loaded in libvips
      2. pdftoppm + pyvips — reliable fallback when pdfload is unavailable
         (pdftoppm renders pages to PNG; pyvips assembles sliding-window WEBP)

    Window structure (pages_per_image pages stacked horizontally):
      Image p001-003: page 1 | page 2 | page 3
      Image p002-004: page 2 | page 3 | page 4
      …
    """
    try:
        _lazy_vips_init(); import pyvips
    except ImportError:
        logger.error("  ✗ pyvips not installed: pip install pyvips")
        return []

    # --- Strategy 1: pyvips pdfload (requires poppler dynamic module) ---
    if _pdfload_available():
        created: List[Path] = []
        try:
            total_pages = get_pdf_page_count(pdf_path)
            logger.info(f"  PDF has {total_pages} page(s)")

            if total_pages <= 0:
                logger.warning(f"  ⚠ No pages found in {pdf_path.name}")
                return []

            for start_page in range(total_pages):
                end_page  = min(start_page + pages_per_image, total_pages)
                page_imgs = []

                for page_num in range(start_page, end_page):
                    page_img = pyvips.Image.pdfload(
                        str(pdf_path),
                        page=page_num,
                        n=1,
                        dpi=IMAGE_DPI,
                        background=[255, 255, 255, 255],
                    )
                    page_imgs.append(page_img)

                img = (
                    pyvips.Image.arrayjoin(page_imgs, across=1)
                    if len(page_imgs) > 1
                    else page_imgs[0]
                )

                out_file = output_dir / f"{base_name}_p{start_page+1:03d}-{end_page:03d}.webp"
                img.webpsave(str(out_file), Q=WEBP_QUALITY, effort=4, smart_subsample=True)
                created.append(out_file)
                logger.info(f"  ✓ Image: pages {start_page+1}-{end_page} → {out_file.name}")

                if end_page >= total_pages:
                    break

            return created

        except Exception as exc:
            logger.warning(f"  ⚠ pdfload failed ({exc}), falling back to pdftoppm…")

    # --- Strategy 2: pdftoppm + pyvips ---
    logger.info("  Using pdftoppm → PNG → WEBP pipeline…")
    return _render_pdf_via_pdftoppm(pdf_path, output_dir, base_name, pages_per_image)

# ---------------------------------------------------------------------------
# Unified image conversion dispatcher
# ---------------------------------------------------------------------------

def convert_to_images(
    filepath: Path,
    output_dir: Path,
    temp_dir: Path,
) -> List[Path]:
    """
    Convert any supported document to WEBP sliding-window images.

    PDF     → render_pdf_pages_to_images (direct pyvips)
    DOCX    → (1) LibreOffice→PDF→pyvips
               (2) docx2pdf→PDF→pyvips
               (3) python-docx native fallback
    PPTX    → (1) LibreOffice→PDF→pyvips
               (2) python-pptx native fallback
    Excel   → (1) LibreOffice→PDF→pyvips
               (2) no images if LibreOffice unavailable (markdown still generated)
    .drawio → draw.io CLI → PDF/PNG → WEBP
    """
    ext       = filepath.suffix.lower()
    base_name = filepath.stem

    # ---- PDF: render directly ----
    if ext == '.pdf':
        return render_pdf_pages_to_images(filepath, output_dir, base_name)

    # ---- DOCX / DOC ----
    if ext in ('.docx', '.doc'):
        logger.info(f"  Converting {ext} → PDF for image extraction…")
        pdf_path = convert_office_to_pdf(filepath, temp_dir)
        if pdf_path and pdf_path.exists():
            return render_pdf_pages_to_images(pdf_path, output_dir, base_name)

        logger.info("  Falling back to python-docx native renderer…")
        return convert_docx_to_images_fallback(filepath, output_dir, base_name)

    # ---- PPTX / PPT ----
    if ext in ('.pptx', '.ppt'):
        logger.info(f"  Converting {ext} → PDF for image extraction…")
        pdf_path = convert_office_to_pdf(filepath, temp_dir)
        if pdf_path and pdf_path.exists():
            return render_pdf_pages_to_images(pdf_path, output_dir, base_name)

        logger.info("  Falling back to python-pptx native renderer…")
        return convert_pptx_to_images_fallback(filepath, output_dir, base_name)

    # ---- Excel: convert via LibreOffice → PDF → WEBP ----
    if ext in ('.xlsx', '.xls', '.xlsm'):
        logger.info(f"  Converting {ext} → PDF for image extraction…")
        pdf_path = convert_office_to_pdf(filepath, temp_dir)
        if pdf_path and pdf_path.exists():
            return render_pdf_pages_to_images(pdf_path, output_dir, base_name)
        logger.warning("  ⚠ Excel→PDF failed — no images generated (markdown still available)")
        return []

    # ---- draw.io diagrams ----
    if ext == '.drawio':
        logger.info("  Converting .drawio → images via draw.io CLI…")
        return convert_drawio_to_images(filepath, output_dir, temp_dir)

    logger.warning(f"  ⚠ No image conversion path for extension: {ext}")
    return []


# ---------------------------------------------------------------------------
# Video processing: VTT subtitles + scene-change cadres
# ---------------------------------------------------------------------------

def _find_companion_subtitles(video_path: Path) -> List[Path]:
    """Find .vtt/.srt files alongside a video in __SPECS__ (manual subtitles)."""
    companions = []
    parent = video_path.parent
    stem = video_path.stem
    for ext in SUBTITLE_EXTENSIONS:
        # Exact match: same stem
        exact = parent / f"{stem}{ext}"
        if exact.exists():
            companions.append(exact)
        # Also look for files that start with the video stem (Teams often adds suffixes)
        for f in parent.glob(f"*{ext}"):
            if f not in companions and f.stem.startswith(stem[:40]):
                companions.append(f)
    return companions


def _copy_companion_subtitles(
    video_path: Path, markdown_dir: Path, fragments_dir: Path,
) -> List[str]:
    """Copy manual VTT/SRT files from __SPECS__ into the fragment's markdown dir."""
    companions = _find_companion_subtitles(video_path)
    copied = []
    for sub in companions:
        dest = markdown_dir / sub.name
        if not dest.exists():
            import shutil as _shutil
            _shutil.copy2(sub, dest)
            logger.info(f"  Copied manual subtitle: {sub.name}")
        copied.append(str(dest.relative_to(fragments_dir)))
    return copied


def generate_vtt_subtitles(
    video_path: Path, output_dir: Path, model_name: str = "base",
) -> Optional[Path]:
    """Transcribe video audio to VTT subtitle file using Whisper."""
    try:
        import whisper
        from whisper.utils import get_writer
    except ImportError:
        logger.error("  openai-whisper not installed — skipping VTT generation")
        return None

    import warnings as _w
    _w.filterwarnings("ignore", category=UserWarning)
    _w.filterwarnings("ignore", category=FutureWarning)

    logger.info(f"  Loading Whisper model '{model_name}'…")
    model = whisper.load_model(model_name)

    logger.info("  Transcribing audio…")
    result = model.transcribe(str(video_path))

    vtt_writer = get_writer("vtt", str(output_dir))
    vtt_writer(result, str(video_path))

    # Whisper names output after the input stem — rename to our convention
    expected = output_dir / f"{video_path.stem}.vtt"
    final = output_dir / f"{video_path.stem}_whisper.vtt"
    if expected.exists() and expected != final:
        expected.rename(final)
    elif not expected.exists() and not final.exists():
        vtts = list(output_dir.glob("*.vtt"))
        if vtts:
            vtts[0].rename(final)
        else:
            logger.warning("  Whisper did not produce a VTT file")
            return None

    logger.info(f"  VTT saved: {final.name}")
    return final


def extract_scene_frames(
    video_path: Path, output_dir: Path, threshold: float = 5.0,
) -> List[Path]:
    """Detect scene changes and save the first frame of each new scene as JPEG."""
    try:
        from scenedetect import detect, ContentDetector
    except ImportError:
        logger.error("  scenedetect not installed — skipping scene extraction")
        return []
    try:
        import cv2
    except ImportError:
        logger.error("  opencv not installed — skipping scene extraction")
        return []

    logger.info(f"  Detecting scene changes (threshold={threshold})…")
    scene_list = detect(str(video_path), ContentDetector(threshold=threshold))
    logger.info(f"  Found {len(scene_list)} scene changes")

    cap = cv2.VideoCapture(str(video_path))
    saved = []

    if not scene_list:
        # Fallback: save first frame
        ok, frame = cap.read()
        if ok:
            p = output_dir / "cadre_000.jpg"
            cv2.imwrite(str(p), frame)
            saved.append(p)
    else:
        for i, scene in enumerate(scene_list):
            cap.set(cv2.CAP_PROP_POS_FRAMES, scene[0].get_frames())
            ok, frame = cap.read()
            if ok:
                p = output_dir / f"cadre_{i:03d}.jpg"
                cv2.imwrite(str(p), frame)
                saved.append(p)

    cap.release()
    logger.info(f"  Saved {len(saved)} scene images")
    return saved


def process_video(
    filepath: Path,
    fragments_dir: Path,
    markdown_dir: Path,
    images_dir: Path,
    display_name: str,
) -> Dict:
    """Process a video file: companion subs + Whisper VTT + scene-change cadres.

    Returns a results dict compatible with the standard manifest format.
    """
    results = {
        "source":            str(filepath),
        "hash":              get_file_hash(filepath),
        "processed_at":      datetime.now().isoformat(),
        "markitdown":        None,
        "docling":           None,
        "drawio_parsed":     None,
        "vtt":               None,
        "manual_subtitles":  [],
        "images":            [],
        "status":            "pending",
        "type":              "video",
    }

    # 1. Copy manual VTT/SRT from __SPECS__ (preserve originals)
    logger.info("→ Checking for manual subtitles alongside video…")
    results["manual_subtitles"] = _copy_companion_subtitles(
        filepath, markdown_dir, fragments_dir,
    )
    if results["manual_subtitles"]:
        logger.info(f"  Found {len(results['manual_subtitles'])} manual subtitle(s)")

    # 2. Whisper-generated VTT (complements manual subs)
    logger.info("→ Generating VTT subtitles (Whisper)…")
    vtt = generate_vtt_subtitles(filepath, markdown_dir, VIDEO_WHISPER_MODEL)
    if vtt:
        results["vtt"] = str(vtt.relative_to(fragments_dir))

    # 3. Scene-change cadre images
    logger.info("→ Extracting scene-change frames…")
    imgs = extract_scene_frames(filepath, images_dir, VIDEO_SCENE_THRESHOLD)
    results["images"] = [str(p.relative_to(fragments_dir)) for p in imgs]

    # Status
    has_text = bool(results["vtt"] or results["manual_subtitles"])
    has_img = bool(results["images"])
    if has_text and has_img:
        results["status"] = "complete"
    elif has_text:
        results["status"] = "vtt_only"
    elif has_img:
        results["status"] = "images_only"
    else:
        results["status"] = "failed"

    return results


# ---------------------------------------------------------------------------
# Single-document processor
# ---------------------------------------------------------------------------

def process_document(
    filepath: Path,
    fragments_dir: Path,
    manifest: Dict,
    force: bool = False,
    source_label: Optional[str] = None,
) -> Dict:
    """
    Process one document: produce dual markdown + WEBP sliding-window images.

    source_label — optional display name (e.g. for files extracted from archives)
    """
    display_name = source_label or filepath.name
    ext = filepath.suffix.lower()

    results = {
        "source":       str(filepath),
        "hash":         get_file_hash(filepath),
        "processed_at": datetime.now().isoformat(),
        "markitdown":   None,
        "docling":      None,
        "drawio_parsed": None,
        "images":       [],
        "status":       "pending",
    }

    # Subdirectory named after the stem of the original/display name
    safe_stem = Path(display_name).stem
    doc_dir   = fragments_dir / safe_stem
    doc_dir.mkdir(parents=True, exist_ok=True)

    images_dir   = doc_dir / "images"
    markdown_dir = doc_dir / "markdown"
    images_dir.mkdir(exist_ok=True)
    markdown_dir.mkdir(exist_ok=True)

    logger.info(f"\n{'='*60}")
    logger.info(f"Processing: {display_name}")
    logger.info(f"{'='*60}")

    # --- Video files get their own pipeline (VTT + scene cadres) ---
    if SUPPORTED_EXTENSIONS.get(ext) == 'video':
        return process_video(
            filepath, fragments_dir, markdown_dir, images_dir, display_name,
        )

    # 1. Markdown conversion — strategy depends on file type
    if ext == '.drawio':
        # .drawio gets its own XML parser (not markitdown/docling)
        logger.info("→ Parsing draw.io XML to Markdown…")
        md_drawio = convert_drawio_to_markdown(filepath, markdown_dir)
        if md_drawio:
            results["drawio_parsed"] = str(md_drawio.relative_to(fragments_dir))
    else:
        # Standard documents get dual markdown conversion
        logger.info("→ Converting to Markdown (markitdown + docling)…")
        md_markitdown = convert_with_markitdown(filepath, markdown_dir)
        if md_markitdown:
            results["markitdown"] = str(md_markitdown.relative_to(fragments_dir))

        md_docling = convert_with_docling(filepath, markdown_dir)
        if md_docling:
            results["docling"] = str(md_docling.relative_to(fragments_dir))

    # 2. WEBP images
    logger.info("→ Converting to WEBP images (sliding window)…")
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        images = convert_to_images(filepath, images_dir, tmp_path)
        results["images"] = [str(img.relative_to(fragments_dir)) for img in images]

    # Status
    has_md  = bool(results["markitdown"] or results["docling"] or results["drawio_parsed"])
    has_img = bool(results["images"])

    if has_md and has_img:
        results["status"] = "complete"
    elif has_md:
        results["status"] = "markdown_only"
    elif has_img:
        results["status"] = "images_only"
    else:
        results["status"] = "failed"

    logger.info(f"→ Status: {results['status']}")
    return results

# ---------------------------------------------------------------------------
# Index builder
# ---------------------------------------------------------------------------

def create_index_file(fragments_dir: Path, manifest: Dict):
    """Write a human-readable INDEX.md listing all processed fragments."""
    index_path = fragments_dir / "INDEX.md"

    with open(index_path, 'w', encoding='utf-8') as fh:
        fh.write("# Document Fragments Index\n\n")
        fh.write(f"*Last updated: {manifest.get('last_updated', 'Unknown')}*\n\n")
        fh.write("## Processed Documents\n\n")

        for filename, info in manifest.get("files", {}).items():
            fh.write(f"### {filename}\n\n")
            fh.write(f"- **Status:** {info.get('status', 'unknown')}\n")
            fh.write(f"- **Source:** {info.get('source', 'unknown')}\n")
            fh.write(f"- **Processed:** {info.get('processed_at', 'unknown')}\n")

            if info.get('markitdown'):
                fh.write(f"- **Markitdown:** [{info['markitdown']}]({info['markitdown']})\n")
            if info.get('docling'):
                fh.write(f"- **Docling:** [{info['docling']}]({info['docling']})\n")
            if info.get('drawio_parsed'):
                fh.write(f"- **Drawio Parsed:** [{info['drawio_parsed']}]({info['drawio_parsed']})\n")

            imgs = info.get('images', [])
            if imgs:
                fh.write(f"- **Images:** {len(imgs)} WEBP file(s)\n")
                for img in imgs[:3]:
                    fh.write(f"  - [{img}]({img})\n")
                if len(imgs) > 3:
                    fh.write(f"  - … and {len(imgs) - 3} more\n")

            if info.get('archive_contents'):
                fh.write(f"- **Archive contents:** {len(info['archive_contents'])} document(s)\n")
                for ac in info['archive_contents']:
                    fh.write(f"  - {ac}\n")

            fh.write("\n")

    logger.info(f"✓ Index → {index_path}")

# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def _compute_display_names(file_list: List[Path]) -> Dict[Path, str]:
    """
    Detect filename collisions among collected files and compute unique
    display names by prepending the nearest distinguishing parent directory.

    When multiple files share the same stem (e.g. files from two different
    subdirectories both named ``report.docx``), the fragment directory would
    collide.  This function returns a mapping  ``{filepath: display_name}``
    where ``display_name`` includes a context prefix for colliding files and
    equals the bare filename for unique ones.

    Prefix derivation (applied per colliding group):
      1. Find the first ancestor directory component that differs among the
         colliding files — that is the distinguishing context.
      2. Clean the component (collapse whitespace/special chars to ``-``).
      3. Prepend as ``{Prefix}-{original_filename}``.
      4. Fallback: MD5 hash of the parent path (guarantees uniqueness).
    """
    from collections import defaultdict

    # Group files by their stem (the part that becomes the fragment dir name)
    stem_groups: Dict[str, List[Path]] = defaultdict(list)
    for fpath in file_list:
        stem_groups[fpath.stem].append(fpath)

    display: Dict[Path, str] = {}

    for stem, paths in stem_groups.items():
        if len(paths) == 1:
            # No collision — use bare filename
            display[paths[0]] = paths[0].name
            continue

        # Collision detected — derive a unique prefix per file
        logger.info(f"⚠ Filename collision on '{stem}' across {len(paths)} files — adding context prefix")

        for fpath in paths:
            prefix = _derive_context_prefix(fpath, paths)
            if prefix:
                unique_name = f"{prefix}-{fpath.name}"
            else:
                # Fallback: use hash of parent path to guarantee uniqueness
                parent_hash = hashlib.md5(str(fpath.parent).encode()).hexdigest()[:6]
                unique_name = f"{parent_hash}-{fpath.name}"
            display[fpath] = unique_name
            logger.info(f"  · {fpath} → fragment '{Path(unique_name).stem}'")

    return display


def _derive_context_prefix(filepath: Path, collision_group: List[Path]) -> Optional[str]:
    """
    Find the first ancestor directory component of ``filepath`` that differs
    from the corresponding component of every other file in ``collision_group``.

    This is purely path-structural — no project-specific keywords.
    Returns a cleaned prefix string, or None if no distinguishing component
    can be found (caller should fall back to a hash).
    """
    # Collect resolved parent parts for each file (excluding the filename itself)
    all_parts = [p.resolve().parts[:-1] for p in collision_group]
    own_parts = filepath.resolve().parts[:-1]

    # Walk own ancestors from closest (immediate parent) outward
    _GENERIC_DIRS = {"meeting notes", "meetings", "documents", "notes", "docs",
                     "templates", "attachments", "files", "shared"}

    for part in reversed(own_parts):
        part_lower = part.lower()
        if part_lower in _GENERIC_DIRS:
            continue
        # Check whether this component is absent or different in at least one
        # other file's path — i.e. it actually distinguishes this file
        others_have_same = all(
            part in other_parts
            for other_parts in all_parts
            if other_parts != own_parts
        )
        if not others_have_same:
            return _clean_prefix(part)

    # All ancestors are shared — use the immediate parent as best effort
    if own_parts:
        return _clean_prefix(own_parts[-1])
    return None


def _clean_prefix(raw: str) -> str:
    """Normalize a directory name into a clean fragment prefix."""
    import re
    # Replace spaces and special chars with hyphens
    cleaned = re.sub(r"[\s/\\]+", "-", raw.strip())
    # Collapse multiple hyphens
    cleaned = re.sub(r"-{2,}", "-", cleaned)
    # Strip leading/trailing hyphens
    cleaned = cleaned.strip("-")
    return cleaned


def _collect_files_recursive(scan_dir: Path) -> tuple:
    """
    Recursively collect supported documents and archives from scan_dir,
    excluding directories in EXCLUDED_DIRS.

    Returns (document_files, archive_files) as two sorted lists of Paths.
    """
    doc_files: List[Path] = []
    arc_files: List[Path] = []

    for dirpath, dirnames, filenames in os.walk(scan_dir, followlinks=True):
        # Prune excluded directories in-place (prevents os.walk from descending)
        dirnames[:] = [
            d for d in dirnames
            if d not in EXCLUDED_DIRS and not d.startswith('.')
        ]

        for fname in filenames:
            fpath = Path(dirpath) / fname
            ext_lower = fpath.suffix.lower()

            # Handle compound extensions like .tar.gz
            if is_archive(fpath):
                arc_files.append(fpath)
            elif ext_lower in SUPPORTED_EXTENSIONS:
                doc_files.append(fpath)

    return sorted(doc_files), sorted(arc_files)


def main():
    parser = argparse.ArgumentParser(
        description="Convert documents to Markdown + WEBP for LLM processing.\n\n"
                    "By default, scans the current directory recursively for all\n"
                    "supported document types. Use --scan-dir to limit the scan\n"
                    "to a specific directory (e.g. __SPECS__/).",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        '--scan-dir', type=Path,
        default=None,
        help="Directory to scan for documents. Default: CWD (recursive). "
             "Use --scan-dir __SPECS__ to limit to that directory.",
    )
    # Keep --specs-dir as an alias for backward compatibility
    parser.add_argument(
        '--specs-dir', type=Path,
        default=None,
        help="(Deprecated alias for --scan-dir) Source directory containing documents.",
    )
    parser.add_argument(
        '--fragments-dir', type=Path,
        default=Path("__FRAGMENTS__"),
        help="Output directory for converted files (default: ./__FRAGMENTS__/ relative to CWD)",
    )
    parser.add_argument(
        '--force', action='store_true',
        help="Force reprocessing even for unchanged files",
    )
    parser.add_argument(
        '--file', type=str,
        help="Process only this specific filename (searched recursively in scan dir)",
    )
    parser.add_argument(
        '--no-recurse', action='store_true',
        help="Only scan top-level of the scan directory (no subdirectories)",
    )
    parser.add_argument(
        '--clean', action='store_true',
        help="Remove fragments for documents that no longer exist in the scan directory",
    )
    parser.add_argument(
        '--whisper-model', type=str, default="base",
        choices=["tiny", "base", "small", "medium", "large"],
        help="Whisper model for video transcription (default: base)",
    )
    parser.add_argument(
        '--scene-threshold', type=float, default=5.0,
        help="Scene detection sensitivity for videos — lower = more cadres (default: 5.0)",
    )

    args = parser.parse_args()

    # Apply video settings to module-level defaults
    global VIDEO_WHISPER_MODEL, VIDEO_SCENE_THRESHOLD
    VIDEO_WHISPER_MODEL = args.whisper_model
    VIDEO_SCENE_THRESHOLD = args.scene_threshold

    # Resolve scan directory: --scan-dir > --specs-dir > auto-detect
    scan_dir = args.scan_dir or args.specs_dir
    if scan_dir is None:
        # Auto-detect: if __SPECS__/ exists, use it; otherwise scan CWD
        if Path("__SPECS__").exists():
            scan_dir = Path("__SPECS__")
            logger.info("Auto-detected __SPECS__/ directory — scanning it.")
        else:
            scan_dir = Path(".")
            logger.info("No __SPECS__/ found — scanning project root recursively.")
    
    if not scan_dir.exists():
        logger.error(f"✗ Scan directory not found: {scan_dir}")
        sys.exit(1)

    args.fragments_dir.mkdir(parents=True, exist_ok=True)

    manifest_path = args.fragments_dir / ".manifest.json"
    manifest      = load_manifest(manifest_path)

    logger.info("=" * 60)
    logger.info("Document Converter Pipeline")
    logger.info("=" * 60)
    logger.info(f"Scan dir:  {scan_dir.resolve()}")
    logger.info(f"Output:    {args.fragments_dir.resolve()}")
    logger.info(f"Recursive: {not args.no_recurse}")
    logger.info(f"Force:     {args.force}")
    if DRAWIO_CLI:
        logger.info(f"draw.io:   {DRAWIO_CLI}")
    else:
        logger.info("draw.io:   not found (diagram images will be skipped)")

    # ------------------------------------------------------------------ #
    # Collect files                                                        #
    # ------------------------------------------------------------------ #
    if args.no_recurse:
        # Flat scan — top-level only (original behaviour)
        all_files: List[Path] = []
        for ext in SUPPORTED_EXTENSIONS:
            all_files.extend(scan_dir.glob(f"*{ext}"))
            all_files.extend(scan_dir.glob(f"*{ext.upper()}"))
        all_files = sorted(set(all_files))

        archive_files: List[Path] = []
        for candidate in sorted(scan_dir.iterdir()):
            if candidate.is_file() and is_archive(candidate):
                archive_files.append(candidate)
    else:
        # Recursive scan with directory exclusions
        all_files, archive_files = _collect_files_recursive(scan_dir)

    # Apply --file filter to documents AND archives
    if args.file:
        target_name = args.file
        matching_docs = [f for f in all_files if f.name == target_name]
        matching_archives = [a for a in archive_files if a.name == target_name]

        if matching_docs:
            # Target is a regular document — skip all archives
            all_files = matching_docs
            archive_files = []
        elif matching_archives:
            # Target is an archive — process only that archive, skip regular docs
            all_files = []
            archive_files = matching_archives
        else:
            logger.error(f"✗ File not found in {scan_dir}: {target_name}")
            sys.exit(1)

    logger.info(f"Found {len(all_files)} document(s), {len(archive_files)} archive(s)")

    # Build the set of manifest keys for all currently-found files.
    # Used for change tracking AND orphan detection.
    current_file_keys: set = set()
    for f in all_files:
        current_file_keys.add(_manifest_key(f))
    for a in archive_files:
        current_file_keys.add(_manifest_key(a))

    prev_count = len(manifest.get("files", {}))
    if prev_count > 0:
        logger.info(f"Manifest: {prev_count} previously tracked entry/entries")

    processed = 0
    skipped   = 0
    failed    = 0

    # ------------------------------------------------------------------ #
    # Process archives — extract then process each contained document     #
    # ------------------------------------------------------------------ #
    for archive_path in sorted(archive_files):
        file_key = _manifest_key(archive_path)
        if not needs_processing(archive_path, manifest, args.force, file_key=file_key):
            logger.info(f"⏭ Skipping archive (unchanged): {archive_path.name}")
            skipped += 1
            continue

        logger.info(f"\n{'='*60}")
        logger.info(f"Archive: {archive_path.name}")
        logger.info(f"{'='*60}")

        with tempfile.TemporaryDirectory() as arc_tmp:
            arc_tmp_path = Path(arc_tmp)
            extracted = extract_archive(archive_path, arc_tmp_path)

            archive_record = {
                "source":           str(archive_path),
                "hash":             get_file_hash(archive_path),
                "processed_at":     datetime.now().isoformat(),
                "type":             "archive",
                "archive_contents": [],
                "status":           "complete" if extracted else "empty",
            }

            for inner_doc in extracted:
                # Use archive_name/doc_name as the display label
                rel = inner_doc.relative_to(arc_tmp_path)
                display = f"{archive_path.stem}/{rel}"

                try:
                    doc_results = process_document(
                        inner_doc,
                        args.fragments_dir,
                        manifest,
                        force=args.force,
                        source_label=display,
                    )
                    inner_key = display
                    manifest["files"][inner_key] = doc_results
                    archive_record["archive_contents"].append(inner_key)

                    if doc_results["status"] in ("complete", "markdown_only", "images_only"):
                        processed += 1
                    else:
                        failed += 1

                except Exception as exc:
                    logger.error(f"✗ Error processing {display}: {exc}")
                    failed += 1

            manifest["files"][file_key] = archive_record

    # ------------------------------------------------------------------ #
    # Compute collision-safe display names for all regular documents       #
    # ------------------------------------------------------------------ #
    display_names = _compute_display_names(all_files)

    # ------------------------------------------------------------------ #
    # Process regular documents                                            #
    # ------------------------------------------------------------------ #
    for doc_path in sorted(all_files):
        file_key = _manifest_key(doc_path)
        if not needs_processing(doc_path, manifest, args.force, file_key=file_key):
            logger.info(f"⏭ Skipping (unchanged): {doc_path.name}")
            skipped += 1
            continue

        # Use collision-safe display name (includes partner prefix when needed)
        label = display_names.get(doc_path)

        try:
            results = process_document(
                doc_path, args.fragments_dir, manifest, args.force,
                source_label=label,
            )
            manifest["files"][file_key] = results

            if results["status"] in ("complete", "markdown_only", "images_only"):
                processed += 1
            else:
                failed += 1

        except Exception as exc:
            logger.error(f"✗ Error processing {doc_path.name}: {exc}")
            manifest["files"][file_key] = {
                "source":       str(doc_path),
                "hash":         get_file_hash(doc_path),
                "processed_at": datetime.now().isoformat(),
                "status":       "error",
                "error":        str(exc),
            }
            failed += 1

        # Persist manifest after every file so an interrupted run does not
        # lose progress — the next invocation will skip already-completed files.
        save_manifest(manifest_path, manifest)

    # ------------------------------------------------------------------ #
    # Orphan detection & cleanup                                           #
    # ------------------------------------------------------------------ #
    orphaned_keys = _detect_orphaned_entries(manifest, current_file_keys)
    cleaned = 0

    if orphaned_keys:
        if args.clean:
            logger.info(f"\n→ Cleaning {len(orphaned_keys)} orphaned entry/entries…")
            cleaned = _clean_orphaned_fragments(orphaned_keys, manifest, args.fragments_dir)
        else:
            logger.info(f"\n⚠ {len(orphaned_keys)} orphaned manifest entry/entries "
                        f"(source files no longer found):")
            for okey in orphaned_keys:
                logger.info(f"  · {okey}")
            logger.info("  Run with --clean to remove their fragments.")

    # ------------------------------------------------------------------ #
    # Persist & report                                                     #
    # ------------------------------------------------------------------ #
    save_manifest(manifest_path, manifest)
    create_index_file(args.fragments_dir, manifest)

    logger.info("\n" + "=" * 60)
    logger.info("SUMMARY")
    logger.info("=" * 60)
    logger.info(f"✓ Processed:  {processed}  (new or changed)")
    logger.info(f"⏭ Unchanged:  {skipped}")
    logger.info(f"✗ Failed:     {failed}")
    if cleaned:
        logger.info(f"🗑 Cleaned:    {cleaned}  (orphaned fragments removed)")
    elif orphaned_keys:
        logger.info(f"⚠ Orphaned:   {len(orphaned_keys)}  (use --clean to remove)")
    logger.info(f"📁 Output:    {args.fragments_dir}")
    logger.info(f"📋 Tracked:   {len(manifest.get('files', {}))} total entry/entries")
    logger.info("=" * 60)


if __name__ == "__main__":
    main()